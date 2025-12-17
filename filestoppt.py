# app.py
import io
import os
import sys
import zipfile
import tempfile
import shutil
import subprocess
import traceback
from pathlib import Path
from flask import Flask, request, make_response, jsonify
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches
import fitz  # PyMuPDF
from PIL import Image, ImageDraw, ImageFont

app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}})

# Supported extensions (lowercase)
SUPPORTED_SINGLE = {
    '.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx',
    '.html', '.htm', '.jpg', '.jpeg', '.png', '.webp', '.txt'
}

# Utilities -------------------------------------------------------------------

def which_bin(name):
    """Return path to executable or None."""
    from shutil import which
    return which(name)

SOFFICE_BIN = which_bin('soffice') or which_bin('libreoffice')  # common names
WKHTMLTOPDF_BIN = which_bin('wkhtmltopdf')

def run_soffice_convert(input_path: str, output_dir: str):
    """
    Convert a document (office/html) to PDF using LibreOffice headless.
    Creates .pdf in output_dir. Returns path to PDF on success.
    """
    if not SOFFICE_BIN:
        raise RuntimeError("LibreOffice (soffice) not found on PATH. Install LibreOffice.")
    # LibreOffice converts to output_dir with --convert-to pdf --outdir
    cmd = [SOFFICE_BIN, '--headless', '--convert-to', 'pdf:writer_pdf_Export', '--outdir', output_dir, input_path]
    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    # find produced pdf in output_dir
    base = os.path.splitext(os.path.basename(input_path))[0]
    pdf_path = os.path.join(output_dir, base + '.pdf')
    if not os.path.exists(pdf_path):
        # Some files may produce different name; try searching for a created .pdf
        for f in os.listdir(output_dir):
            if f.lower().endswith('.pdf') and os.path.splitext(f)[0].lower().startswith(base.lower()):
                return os.path.join(output_dir, f)
        raise RuntimeError(f"LibreOffice did not produce PDF for {input_path}")
    return pdf_path

def run_wkhtmltopdf(input_path: str, output_pdf: str):
    """Use wkhtmltopdf if installed (for HTML -> PDF)."""
    if not WKHTMLTOPDF_BIN:
        raise RuntimeError("wkhtmltopdf not found")
    cmd = [WKHTMLTOPDF_BIN, input_path, output_pdf]
    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    if not os.path.exists(output_pdf):
        raise RuntimeError("wkhtmltopdf did not produce output")

# Conversion helpers ----------------------------------------------------------

def pdf_bytes_to_pptx_bytes(pdf_bytes: bytes):
    """
    Convert PDF bytes to a PPTX bytes.
    Each page becomes an image slide.
    """
    doc = fitz.open(stream=pdf_bytes, filetype='pdf')
    prs = Presentation()
    # for each page rasterize to PNG and add as picture to slide
    for i in range(doc.page_count):
        page = doc.load_page(i)
        # choose zoom for reasonable quality (tweak if needed)
        zoom = 2  # 2 -> 144 DPI (72*2)
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img_bytes = pix.tobytes(output='png')
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        img_stream = io.BytesIO(img_bytes)
        # measure image size in pixels
        img = Image.open(io.BytesIO(img_bytes))
        px_w, px_h = img.size
        dpi = 96.0
        width_in = px_w / dpi
        height_in = px_h / dpi
        prs_w_in = prs.slide_width / 914400.0
        prs_h_in = prs.slide_height / 914400.0
        scale = min(prs_w_in / width_in, prs_h_in / height_in)
        pic_w_in = width_in * scale
        pic_h_in = height_in * scale
        left_in = (prs_w_in - pic_w_in) / 2
        top_in = (prs_h_in - pic_h_in) / 2
        slide.shapes.add_picture(img_stream, Inches(left_in), Inches(top_in), width=Inches(pic_w_in), height=Inches(pic_h_in))
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    doc.close()
    return out.read()

def image_bytes_to_pptx_bytes(img_bytes: bytes, img_name: str = 'image'):
    """Create a PPTX with the image as a single slide."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    img = Image.open(io.BytesIO(img_bytes))
    px_w, px_h = img.size
    dpi = 96.0
    width_in = px_w / dpi
    height_in = px_h / dpi
    prs_w_in = prs.slide_width / 914400.0
    prs_h_in = prs.slide_height / 914400.0
    scale = min(prs_w_in / width_in, prs_h_in / height_in)
    pic_w_in = width_in * scale
    pic_h_in = height_in * scale
    left_in = (prs_w_in - pic_w_in) / 2
    top_in = (prs_h_in - pic_h_in) / 2
    slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(left_in), Inches(top_in), width=Inches(pic_w_in), height=Inches(pic_h_in))
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()

def text_to_image_bytes(text: str, max_width=1200, font_path=None, font_size=18):
    """Render text to an image and return PNG bytes."""
    # Choose a font (Pillow default or load a TTF if available)
    try:
        if font_path and os.path.exists(font_path):
            font = ImageFont.truetype(font_path, font_size)
        else:
            font = ImageFont.load_default()
    except Exception:
        font = ImageFont.load_default()
    # simple wrap: measure lines
    lines = []
    words = text.split()
    cur = ""
    # approximate wrapping by characters (simple)
    approx_chars_per_line = 80
    for w in words:
        if len(cur) + 1 + len(w) > approx_chars_per_line:
            lines.append(cur)
            cur = w
        else:
            if cur:
                cur += ' ' + w
            else:
                cur = w
    if cur:
        lines.append(cur)
    # fallback if no words (empty file)
    if not lines:
        lines = ['']

    # compute image size
    line_height = font_size + 6
    img_h = max(300, line_height * len(lines) + 40)
    img_w = max_width
    image = Image.new('RGB', (img_w, img_h), color='white')
    draw = ImageDraw.Draw(image)
    y = 20
    x = 20
    for line in lines:
        draw.text((x, y), line, fill='black', font=font)
        y += line_height
    buf = io.BytesIO()
    image.save(buf, format='PNG')
    buf.seek(0)
    return buf.read()

# High-level conversion ------------------------------------------------------

def convert_supported_file_to_pptx_bytes(file_path: str, ext: str, tmpdir: str):
    """
    Convert the file at file_path to PPTX bytes.
    ext is lowercase extension including dot.
    tmpdir is a working folder for temporary outputs.
    """
    try:
        if ext == '.pdf':
            with open(file_path, 'rb') as f:
                pdf_bytes = f.read()
            return pdf_bytes_to_pptx_bytes(pdf_bytes)

        if ext in {'.jpg', '.jpeg', '.png', '.webp'}:
            with open(file_path, 'rb') as f:
                img_bytes = f.read()
            return image_bytes_to_pptx_bytes(img_bytes, img_name=os.path.basename(file_path))

        if ext == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                txt = f.read()
            img_b = text_to_image_bytes(txt)
            return image_bytes_to_pptx_bytes(img_b, img_name=os.path.basename(file_path))

        if ext in {'.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.html', '.htm'}:
            # Convert to PDF first using soffice (LibreOffice) or wkhtmltopdf for HTML if available
            if ext in {'.html', '.htm'} and WKHTMLTOPDF_BIN:
                out_pdf = os.path.join(tmpdir, os.path.splitext(os.path.basename(file_path))[0] + '.pdf')
                try:
                    run_wkhtmltopdf(file_path, out_pdf)
                    with open(out_pdf, 'rb') as f:
                        pdf_bytes = f.read()
                    return pdf_bytes_to_pptx_bytes(pdf_bytes)
                except Exception:
                    # fallback to libreoffice path below
                    pass
            # Use LibreOffice conversion
            pdf_path = run_soffice_convert(file_path, tmpdir)
            with open(pdf_path, 'rb') as f:
                pdf_bytes = f.read()
            return pdf_bytes_to_pptx_bytes(pdf_bytes)

        raise RuntimeError(f"Unsupported extension: {ext}")
    except Exception as e:
        # rethrow with stacktrace for logging
        raise

# Flask endpoint -------------------------------------------------------------

@app.route('/convert-all-to-ppt', methods=['POST'])
def convert_all_to_ppt():
    """
    Accepts 'file' multipart form.
    If single file -> return single PPTX.
    If ZIP -> convert each supported file inside and return a ZIP of results.
    """
    if 'file' not in request.files:
        return jsonify({'error': 'No file field'}), 400
    uploaded = request.files['file']
    filename = uploaded.filename or 'upload'
    ext = Path(filename).suffix.lower()

    # Working temporary directory (unique)
    base_tmp = tempfile.mkdtemp(prefix='convert_all_')
    try:
        if ext == '.zip':
            zip_path = os.path.join(base_tmp, 'uploaded.zip')
            uploaded.save(zip_path)
            results = []  # list of (name, bytes)
            with zipfile.ZipFile(zip_path, 'r') as zf:
                namelist = zf.namelist()
                for member in namelist:
                    # skip directories
                    if member.endswith('/'):
                        continue
                    member_ext = Path(member).suffix.lower()
                    if member_ext not in SUPPORTED_SINGLE:
                        # skip unsupported — but could add as notice
                        continue
                    # extract member to temp file
                    safe_name = os.path.basename(member)
                    out_path = os.path.join(base_tmp, safe_name)
                    # ensure target directory exists
                    os.makedirs(os.path.dirname(out_path), exist_ok=True)
                    with open(out_path, 'wb') as out_f:
                        out_f.write(zf.read(member))
                    try:
                        pptx_bytes = convert_supported_file_to_pptx_bytes(out_path, member_ext, base_tmp)
                        out_ppt_name = os.path.splitext(safe_name)[0] + '.pptx'
                        results.append((out_ppt_name, pptx_bytes))
                    except Exception as e:
                        # on error produce a small txt explaining failure
                        tb = traceback.format_exc()
                        err_name = os.path.splitext(safe_name)[0] + '_error.txt'
                        results.append((err_name, f'Failed to convert {member}: {str(e)}\n\n{tb}'.encode('utf-8')))

            if not results:
                return jsonify({'error': 'No supported files found inside ZIP.'}), 400

            # If only one pptx result, return it directly
            if len(results) == 1 and results[0][0].lower().endswith('.pptx'):
                name, data = results[0]
                resp = make_response(data)
                resp.headers.set('Content-Type', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
                resp.headers.set('Content-Disposition', f'attachment; filename="{name}"')
                return resp

            # else create zip
            out_zip = io.BytesIO()
            with zipfile.ZipFile(out_zip, 'w', compression=zipfile.ZIP_DEFLATED) as zout:
                for name, data in results:
                    zout.writestr(name, data)
            out_zip.seek(0)
            resp = make_response(out_zip.read())
            resp.headers.set('Content-Type', 'application/zip')
            resp.headers.set('Content-Disposition', 'attachment; filename="converted_ppts.zip"')
            return resp

        else:
            # single file
            save_path = os.path.join(base_tmp, filename)
            uploaded.save(save_path)
            if ext not in SUPPORTED_SINGLE:
                return jsonify({'error': f'Unsupported file type: {ext}'}), 400
            try:
                pptx_bytes = convert_supported_file_to_pptx_bytes(save_path, ext, base_tmp)
                out_name = os.path.splitext(filename)[0] + '.pptx'
                resp = make_response(pptx_bytes)
                resp.headers.set('Content-Type', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
                resp.headers.set('Content-Disposition', f'attachment; filename="{out_name}"')
                return resp
            except Exception as e:
                tb = traceback.format_exc()
                return jsonify({'error': f'Conversion failed: {str(e)}', 'trace': tb}), 500
    finally:
        # cleanup
        try:
            shutil.rmtree(base_tmp)
        except Exception:
            pass

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
